from pptx import Presentation
from pptx.util import Emu
import os, glob

prs = Presentation()
prs.slide_width = Emu(9144000 * 16 / 9 / 16 * 16)
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)

blank = prs.slide_layouts[6]

png_dir = os.path.join(os.path.dirname(__file__), 'png')
files = sorted(glob.glob(os.path.join(png_dir, 'slide*.png')))

for f in files:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(f, 0, 0, width=prs.slide_width, height=prs.slide_height)

out = os.path.join(os.path.dirname(__file__), 'Webinar Solar v2 - Demanda + IA + Eletropostos - Slides.pptx')
prs.save(out)
print(out)
